"""
Generates Final Project Presentation (.pptx) for Kyamatu SMS
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT = os.path.dirname(os.path.abspath(__file__))
DARK = RGBColor(0x1B, 0x1B, 0x2F)
ACCENT = RGBColor(0x00, 0x7B, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xE0, 0xE0, 0xE0)
GREEN = RGBColor(0x00, 0xC8, 0x53)

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf

def add_bullet_frame(slide, left, top, width, height, items, size=16, color=WHITE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return tf

def title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, DARK)
    add_text(slide, 0.5, 1.5, 9, 1.5, title, size=32, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 0.5, 3.2, 9, 1, subtitle, size=18, color=LIGHT, align=PP_ALIGN.CENTER)
    return slide

def section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_text(slide, 0.5, 2.5, 9, 1, title, size=36, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    return slide

def content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_text(slide, 0.5, 0.3, 9, 0.8, title, size=28, color=ACCENT, bold=True)
    # Divider line
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, Inches(0.5), Inches(1.05), Inches(9), Inches(0.02))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    add_bullet_frame(slide, 0.6, 1.2, 8.5, 5, bullets, size=16, color=WHITE)
    return slide

def two_col_slide(prs, title, left_items, right_items, left_title="", right_title=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_text(slide, 0.5, 0.3, 9, 0.8, title, size=28, color=ACCENT, bold=True)
    shape = slide.shapes.add_shape(1, Inches(0.5), Inches(1.05), Inches(9), Inches(0.02))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    if left_title:
        add_text(slide, 0.5, 1.15, 4.3, 0.5, left_title, size=18, color=GREEN, bold=True)
    if right_title:
        add_text(slide, 5.2, 1.15, 4.3, 0.5, right_title, size=18, color=GREEN, bold=True)
    y_off = 1.7 if left_title else 1.2
    add_bullet_frame(slide, 0.6, y_off, 4.2, 5, left_items, size=14, color=WHITE)
    add_bullet_frame(slide, 5.3, y_off, 4.2, 5, right_items, size=14, color=WHITE)
    return slide

def table_slide(prs, title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    add_text(slide, 0.5, 0.3, 9, 0.8, title, size=28, color=ACCENT, bold=True)
    cols = len(headers)
    tbl_rows = len(rows) + 1
    tbl = slide.shapes.add_table(tbl_rows, cols, Inches(0.5), Inches(1.2), Inches(9), Inches(0.4 * tbl_rows)).table
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri+1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x2A, 0x2A, 0x40) if ri % 2 == 0 else RGBColor(0x22, 0x22, 0x38)
    return slide

# ── BUILD PRESENTATION ──────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# 1. Title
s = title_slide(prs, "SCHOOL MANAGEMENT SYSTEM\nFOR KYAMATU PRIMARY SCHOOL", "")
add_text(s, 0.5, 4.3, 9, 0.5, "Mwange Musa Muthami  •  23/05037", size=16, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.5, 4.9, 9, 0.5, "Supervisor: Dr. Simon N. Mwendia / Dr. Kevin Mugoye Sindu", size=14, color=LIGHT, align=PP_ALIGN.CENTER)
add_text(s, 0.5, 5.5, 9, 0.5, "Bachelor of Information Technology  •  BIT 3105  •  April 2026", size=14, color=LIGHT, align=PP_ALIGN.CENTER)

# 2. Problem Statement
content_slide(prs, "Problem Statement", [
    "Kyamatu Primary School relies on paper-based records for all operations",
    "30% of staff time spent on manual transcription of attendance and grades",
    "Physical records vulnerable to loss, damage, and unauthorized alteration",
    "15% data inconsistency rate across term transitions (MoE, 2023)",
    "Parents have no remote access to student academic or financial data",
    "No automated CBC competency grading — everything computed by hand",
])

# 3. Objectives
two_col_slide(prs, "Project Objectives",
    ["Centralized cloud platform for all school data",
     "Automated CBC grading and report card generation",
     "Digital attendance tracking with statistics",
     "Fee management with payment tracking"],
    ["Investigate data loss in current manual systems",
     "Evaluate digital literacy of stakeholders",
     "Design scalable CBC-aligned data model",
     "Implement and test full-stack web application"],
    "System Objectives", "Research Objectives")

# 4. Tech Stack
table_slide(prs, "Technology Stack",
    ["Layer", "Technology", "Purpose"],
    [["Frontend", "React 18 + Vite + Tailwind CSS", "Single Page Application UI"],
     ["Backend", "Node.js 20 + Express.js", "RESTful API Server"],
     ["Database", "PostgreSQL 15 + Prisma ORM", "Relational Data Storage"],
     ["Auth", "JWT + bcrypt (12 rounds)", "Authentication & Security"],
     ["State", "Zustand", "Frontend State Management"],
     ["Hosting", "Cloudflare Pages + Render + Supabase", "Cloud Deployment"]])

# 5. Architecture
content_slide(prs, "System Architecture", [
    "Three-Tier Client-Server Architecture",
    "Client Tier: React 18 SPA served via Cloudflare Pages",
    "Server Tier: Express.js API with 9-layer middleware pipeline",
    "Database Tier: PostgreSQL with 23+ tables, UUID keys, indexed columns",
    "Middleware: CORS → Helmet → Rate Limit → Auth → RBAC → Validate",
    "Auto-deploy via GitHub to Render (backend) and Cloudflare (frontend)",
])

# 6. Database
content_slide(prs, "Database Design", [
    "23+ tables designed with Prisma ORM on PostgreSQL",
    "UUID-based primary keys for all entities",
    "Key entities: User, Student, Staff, Class, Grade, Subject, Attendance",
    "Assessment, AssessmentScore, ReportCard, FeeStructure, Invoice, Payment",
    "Unique constraints prevent duplicate records (e.g., attendance per day)",
    "Indexed columns for fast lookups: email, admission number, dates",
    "6 enums: Role, AdmissionStatus, AttendanceStatus, CompetencyRating, PaymentMethod, PaymentStatus",
])

# 7. Modules
two_col_slide(prs, "Modules Implemented (11 Modules — 100% Complete)",
    ["✅ Authentication & RBAC (6 roles)",
     "✅ Student Management (CRUD + admissions)",
     "✅ Staff Management (assignments)",
     "✅ Academic Structure (CBC grades, terms)",
     "✅ Attendance Tracking (daily marking)",
     "✅ Assessment & Grading (CBC + traditional)"],
    ["✅ Report Cards (rankings, averages)",
     "✅ Fees Management (invoices, payments)",
     "✅ Communication (announcements, DMs)",
     "✅ Timetable (weekly grid)",
     "✅ Dashboards (Admin + Student)",
     "✅ Deployed to Production"])

# 8. Security
content_slide(prs, "Security Features", [
    "JWT authentication with 15-minute access tokens + 7-day refresh tokens",
    "Password hashing with bcrypt (12 salt rounds)",
    "Role-Based Access Control (RBAC) middleware on every route",
    "Redis-backed rate limiting to prevent brute-force attacks",
    "Helmet.js for XSS, CSRF, and clickjacking protection",
    "Input validation (express-validator) + sanitization middleware",
    "Comprehensive audit log (action, entity, user, IP, timestamp)",
])

# 9. Frontend Pages
table_slide(prs, "Frontend — 17 Pages Implemented",
    ["Page", "Description", "Access"],
    [["Login", "Email/password authentication", "All"],
     ["Dashboard", "Role-specific analytics & widgets", "All"],
     ["Students", "Student list with search, filter, CRUD", "Admin/Teacher"],
     ["Admissions", "Approval/rejection workflow", "Admin"],
     ["Staff", "Staff management & assignments", "Admin"],
     ["Attendance", "Daily marking interface", "Teacher"],
     ["Assessments", "Score entry with auto-grading", "Teacher"],
     ["Reports", "Report card generation", "Admin/Teacher"],
     ["Fees", "Invoices, payments, balances", "Bursar"],
     ["Timetable", "Weekly grid view", "All"]])

# 10. Testing
table_slide(prs, "Testing Summary",
    ["Test Type", "Tool", "Tests", "Pass Rate"],
    [["Unit Testing", "Vitest", "Key service functions", "100%"],
     ["Integration Testing", "Postman", "All API endpoints", "100%"],
     ["User Acceptance", "Manual", "End-to-end workflows", "100%"],
     ["Cross-Browser", "Manual", "Chrome, Firefox, Edge, Safari", "100%"],
     ["Responsive", "DevTools", "320px – 1920px viewports", "100%"],
     ["Security", "Manual + Helmet", "Auth, RBAC, rate limiting", "100%"],
     ["Total", "", "46 test cases", "100%"]])

# 11. Challenges
table_slide(prs, "Challenges & Solutions",
    ["Challenge", "Solution"],
    [["Supabase connection failures", "Connection pooling + Prisma retry logic"],
     ["Data loss after page refresh", "Removed auto-seeding from login flow"],
     ["Assessments not grouped by term", "Refactored query with term relations"],
     ["Dashboard showing placeholder data", "Replaced with live API calls"],
     ["Rate limiting in production", "Redis-backed store with graceful fallback"]])

# 12. Demo
section_slide(prs, "🖥️  LIVE DEMO")

# 13. Future Work
content_slide(prs, "Future Enhancements", [
    "M-Pesa STK Push Integration — direct mobile money payments",
    "SMS / Email Notifications — automated alerts for attendance & fees",
    "Flutter Parent App — native mobile app for parents",
    "PDF Report Cards — server-side PDF generation with PDFKit",
    "Multi-School Support — extend architecture for multiple institutions",
])

# 14. Conclusion
content_slide(prs, "Conclusion", [
    "All 11 modules successfully designed, developed, tested, and deployed",
    "100% functional requirement coverage achieved",
    "System addresses: manual records, data loss, slow reports, poor communication",
    "Fully compliant with Kenya's CBC grading framework",
    "User acceptance testing completed with school administrators",
    "System is live and accessible to all stakeholders",
])

# 15. Thank You
s = title_slide(prs, "THANK YOU", "Questions?")
add_text(s, 0.5, 4.5, 9, 0.5, "Mwange Musa Muthami  •  23/05037", size=16, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.5, 5.1, 9, 0.5, "GitHub: https://github.com/MuthamiM/KYAMATU-SMS", size=14, color=LIGHT, align=PP_ALIGN.CENTER)

# Save
path = os.path.join(OUT, "8_Final_Presentation.pptx")
prs.save(path)
print(f"✓ 8_Final_Presentation.pptx ({len(prs.slides)} slides)")
